#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Construit le PowerPoint de revue formative PEV — Antenne Kinshasa Est.

- Diapo 10 : Résultats de processus (Jan-Avr 2025 vs Jan-Avr 2026) — table existante remplie.
- Diapo 11 : Performances par AS et ZS (Jan-Avr 2026) — complétude/promptitude, CV admin,
  enfants ZD/SV, catégorisation ACZ, disponibilité des intrants, taux de perte.
- Diapo 12 : Performance communication — enfants identifiés/récupérés (RECO),
  taux d'abandon BCG-VAR1 et DTC1-DTC3.

Sources (toutes dans le repo) :
  docs/data_as/dashboard/by_as/kn_kinshasa_province.json.gz   (AS, avec population)
  docs/data_as/monthly/<YM>/*.ndjson + docs/data_as/ou_map_as.json + antenne_rules.json
Formules : strictement alignées sur le dashboard (docs/index.html).
"""
import gzip, json, glob, re, copy, os
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PROV_FILE = os.path.join(ROOT, "docs/data_as/dashboard/by_as/kn_kinshasa_province.json.gz")
OU_MAP = os.path.join(ROOT, "docs/data_as/ou_map_as.json")
ANT_RULES = os.path.join(ROOT, "docs/config/antenne_rules.json")
TEMPLATE = os.path.join(ROOT, "public/canevas/canevas_revue_formative_pev.pptx")
OUT = os.path.join(ROOT, "Canevas_Revue_Formative_PEV_Kinshasa_Est.pptx")

ANTENNE = "Kin Est"
M2026 = ["202601", "202602", "202603", "202604"]
M2025 = ["202501", "202502", "202503", "202504"]
NP = 4
NV_RATE, NS_RATE = 0.04, 0.0349
PER_LABEL_26 = "Janvier – Avril 2026"

# ---------- helpers ----------
def num(v):
    try:
        f = float(v)
        return f if f == f else 0.0
    except (TypeError, ValueError):
        return 0.0

def clean(s):
    s = re.sub(r'^[a-z]{2}\s+', '', str(s))
    s = re.sub(r'\s+(Province|Zone de Santé|Aire de Santé|Aire de santé)$', '', s, flags=re.I)
    return s.strip()

def s(recs, f):
    return sum(num(r.get(f)) for r in recs)

# ---------- load ----------
data = json.load(gzip.open(PROV_FILE))
est_all = [r for r in data if r.get("_Antenne") == ANTENNE]
cur = [r for r in est_all if r["_YM"] in M2026]
prev = [r for r in est_all if r["_YM"] in M2025]

# ---------- targets / CV ----------
def targets(recs):
    """Cible NV/NS sur la période : pop (dédupliquée par AS) × taux × NP/12."""
    seen, pop = set(), 0.0
    for r in recs:
        k = (r.get("_ZS"), r.get("_AS"))
        if k not in seen:
            seen.add(k); pop += num(r.get("Pop_par_AS"))
    return {"pop": pop, "nv": pop * NV_RATE * NP / 12.0, "ns": pop * NS_RATE * NP / 12.0}

# Antigènes affichés (ordre calendrier vaccinal) + champ doses + dénominateur
AG_FIELD = {
    "BCG": "BCG_0_11", "VPO0": "VPO0_0_11", "DTC1": "DTC1_0_11", "VPO1": "VPO1_0_11",
    "ROTA1": "ROTA1_0_11", "PCV13-1": "PCV13_1_0_11", "DTC2": "DTC2_0_11", "VPO2": "VPO2_0_11",
    "ROTA2": "ROTA2_0_11", "PCV13-2": "PCV13_2_0_11", "DTC3": "DTC3_0_11", "VPO3": "VPO3_0_11",
    "VPI1": "VPI1_0_11", "ROTA3": "ROTA3_0_11", "PCV13-3": "PCV13_3_0_11", "VAR1": "VAR1_0_11",
    "VPI2": "VPI2_0_11", "VAA": "VAA_0_11", "VAR2": "VAR2_12_23", "Td2+": "Td_2_plus",
}
AG_ORDER = list(AG_FIELD.keys())
AG_USE_NV = {"BCG", "Td2+"}

def cv_pct(recs, ag, tgt):
    den = tgt["nv"] if ag in AG_USE_NV else tgt["ns"]
    if den <= 0:
        return None
    return s(recs, AG_FIELD[ag]) / den * 100.0

def abandon(a, b):
    a, b = num(a), num(b)
    return None if a <= 0 else (a - b) / a * 100.0

def categorize(cv1, ta):
    if cv1 is None or ta is None or ta < 0:
        return "NA"
    h, l = cv1 >= 90, ta <= 10
    return "Cat. 1" if (h and l) else "Cat. 2" if (h and not l) else "Cat. 3" if (not h and l) else "Cat. 4"

def zd(recs):
    return sum(max(0.0, num(r.get("Pop_par_AS")) * NS_RATE / 12.0 - num(r.get("DTC1_0_11"))) for r in recs)

def sv(recs):
    return sum(max(0.0, num(r.get("Pop_par_AS")) * NS_RATE / 12.0 - num(r.get("DTC3_0_11"))) for r in recs)

# Disponibilité : (Σ jours non-rupture) / (Σ 30 sur mois présents) — aligné dashboard.
def dispo(recs, pfx):
    days, nonrup = 0.0, 0.0
    for r in recs:
        rcv = num(r.get(pfx + "_re_ues")); sf = num(r.get(pfx + "_stock_fin"))
        cmm = num(r.get(pfx + "_cmm")); rup = max(0.0, num(r.get(pfx + "_jours_rupture")))
        if rcv > 0 or sf > 0 or cmm > 0 or rup > 0:
            days += 30.0; nonrup += max(0.0, 30.0 - rup)
    return None if days <= 0 else nonrup / days * 100.0

# Taux de perte : pertes / (administrées + pertes), avec repli utilisées-administrées.
def perte(recs, pfx):
    adm = s(recs, pfx + "_administr_es"); util = s(recs, pfx + "_utilis_es"); pdu = s(recs, pfx + "_pertes")
    if pdu <= 0 and util > adm > 0:
        pdu = util - adm
    den = adm + pdu
    return None if den <= 0 else pdu / den * 100.0

# ---------- groupements ----------
def group(recs, by_zs):
    g = {}
    for r in recs:
        k = clean(r["_ZS"]) if by_zs else (clean(r["_ZS"]), clean(r["_AS"]))
        g.setdefault(k, []).append(r)
    return g

g_as = group(cur, False)
g_zs = group(cur, True)

# Étiquettes AS (désambiguïsation des doublons de noms entre ZS)
as_keys = sorted(g_as.keys(), key=lambda k: (k[0], k[1]))
from collections import Counter
name_count = Counter(k[1] for k in as_keys)
def as_label(k):
    return f"{k[1]} ({k[0]})" if name_count[k[1]] > 1 else k[1]
zs_keys = sorted(g_zs.keys())

# ---------- perdues de vue (RECO) depuis le brut mensuel ----------
oum = json.load(open(OU_MAP))
ar = json.load(open(ANT_RULES))
kin_key = [k for k in ar if "inshasa" in k.lower()][0]
zs2ant = ar[kin_key]
ID_F = ["Perdues_de_vue_identifi_s_Penta1_0_11mois", "Perdues_de_vue_identifi_s_Penta1_12_23mois",
        "Perdues_de_vue_identifi_s_Penta3_0_11mois", "Perdues_de_vue_identifi_s_Penta3_12_23mois",
        "Perdues_de_vue_identifi_s_Penta3_24_59mois", "Perdues_de_vue_identifi_s_24_59mois"]
REC_F = ["Perdues_de_vue_r_cup_r_s_Penta1_0_11mois", "Perdues_de_vue_r_cup_r_s_Penta1_12_23mois",
         "Perdues_de_vue_r_cup_r_s_Penta1_24_59mois", "Perdues_de_vue_r_cup_r_s_Penta3_0_11mois",
         "Perdues_de_vue_r_cup_r_s_Penta3_12_23mois", "Perdues_de_vue_r_cup_r_s_Penta3_24_59mois"]
VAR_REC = "Enfants_r_cup_r_s_VAR1_BCU_FAE"

pdv_as = {}   # (zs, as) -> {id, rec, var}
pdv_zs = {}   # zs -> {...}
for ym in M2026:
    for fn in glob.glob(os.path.join(ROOT, f"docs/data_as/monthly/{ym}/*.ndjson")):
        with open(fn) as fh:
            for line in fh:
                line = line.strip()
                if not line:
                    continue
                r = json.loads(line)
                m = oum.get(r.get("OrgUnit"))
                if not m:
                    continue
                zsc = clean(m.get("Org3", ""))
                if zs2ant.get(zsc) != ANTENNE:
                    continue
                asc = clean(m.get("Org4", ""))
                vid = sum(num(r.get(f)) for f in ID_F)
                vrec = sum(num(r.get(f)) for f in REC_F)
                vvar = num(r.get(VAR_REC))
                for d, k in ((pdv_as, (zsc, asc)), (pdv_zs, zsc)):
                    e = d.setdefault(k, {"id": 0.0, "rec": 0.0, "var": 0.0})
                    e["id"] += vid; e["rec"] += vrec; e["var"] += vvar

# ============================================================
#  CONSTRUCTION PPTX
# ============================================================
prs = Presentation(TEMPLATE)
EMU_W = prs.slide_width
EMU_H = prs.slide_height
# layout vierge : celui avec le moins de placeholders
blank = min(prs.slide_layouts, key=lambda L: len(L.placeholders))

NAVY = RGBColor(0x0B, 0x3A, 0x6F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x29, 0x3B)
BAND = RGBColor(0xEE, 0xF4, 0xFC)

def cv_color(p):
    if p is None: return None
    if p > 100: return RGBColor(0x25, 0x63, 0xEB)
    if p >= 90: return RGBColor(0x16, 0xA3, 0x4A)
    if p >= 60: return RGBColor(0xEA, 0xB3, 0x08)
    return RGBColor(0xDC, 0x26, 0x26)

def ab_color(p):
    if p is None: return None
    if p < 0: return RGBColor(0x7F, 0x1D, 0x1D)
    if p <= 10: return RGBColor(0x16, 0xA3, 0x4A)
    if p < 20: return RGBColor(0xEA, 0xB3, 0x08)
    return RGBColor(0xEF, 0x44, 0x44)

def cat_color(c):
    return {"Cat. 1": RGBColor(0x16, 0xA3, 0x4A), "Cat. 2": RGBColor(0x25, 0x63, 0xEB),
            "Cat. 3": RGBColor(0xEA, 0xB3, 0x08), "Cat. 4": RGBColor(0xDC, 0x26, 0x26),
            "NA": RGBColor(0x7F, 0x1D, 0x1D)}.get(c)

def disp_color(p):
    if p is None: return None
    if p >= 90: return RGBColor(0x16, 0xA3, 0x4A)
    if p >= 60: return RGBColor(0xEA, 0xB3, 0x08)
    return RGBColor(0xDC, 0x26, 0x26)

def perte_color(p):
    if p is None: return None
    if p <= 5: return RGBColor(0x16, 0xA3, 0x4A)
    if p <= 10: return RGBColor(0xEA, 0xB3, 0x08)
    return RGBColor(0xDC, 0x26, 0x26)

def add_slide():
    sl = prs.slides.add_slide(blank)
    for ph in list(sl.placeholders):
        ph._element.getparent().remove(ph._element)
    return sl

def add_title(sl, title, subtitle=None):
    tb = sl.shapes.add_textbox(Cm(0.6), Cm(0.3), EMU_W - Cm(1.2), Cm(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(11); r2.font.italic = True; r2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    return tb

def style_cell(cell, text, *, bold=False, fill=None, color=None, align=PP_ALIGN.LEFT, size=8):
    cell.margin_left = Cm(0.08); cell.margin_right = Cm(0.08)
    cell.margin_top = Cm(0.02); cell.margin_bottom = Cm(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid(); cell.fill.fore_color.rgb = fill
    else:
        cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    f = run.font; f.size = Pt(size); f.bold = bold
    f.color.rgb = color if color is not None else (WHITE if fill is not None and bold else DARK)

TOP_DFLT = Cm(2.3)        # haut de zone contenu (sous le titre)
BOT_MARGIN = Cm(0.6)

def rows_per_page(row_cm, header_cm=0.85, top=TOP_DFLT):
    """Nb max de lignes de DONNÉES qui tiennent verticalement sans débordement."""
    avail_cm = (EMU_H - top - BOT_MARGIN) / 360000.0
    return max(4, int((avail_cm - header_cm) / row_cm))

def add_table(sl, header, rows, *, top=TOP_DFLT, size=8, col_w=None, cell_fills=None,
              first_align=PP_ALIGN.LEFT, row_cm=0.6, header_cm=0.85):
    """Tableau à hauteurs de lignes contrôlées (jamais de débordement)."""
    nrows = len(rows) + 1
    ncols = len(header)
    left = Cm(0.5); width = EMU_W - Cm(1.0)
    height = int(Cm(header_cm) + Cm(row_cm) * len(rows))
    gf = sl.shapes.add_table(nrows, ncols, left, top, width, height)
    tbl = gf.table
    tbl.first_row = False; tbl.horz_banding = False
    # hauteurs de lignes explicites
    tbl.rows[0].height = int(Cm(header_cm))
    for r in range(1, nrows):
        tbl.rows[r].height = int(Cm(row_cm))
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = int(width * w / total)
    for c, h in enumerate(header):
        style_cell(tbl.cell(0, c), h, bold=True, fill=NAVY, color=WHITE,
                   align=PP_ALIGN.CENTER, size=size)
    for ri, row in enumerate(rows):
        band = BAND if ri % 2 else WHITE
        for ci, val in enumerate(row):
            fill = band; color = DARK
            if cell_fills and cell_fills[ri][ci] is not None:
                fill = cell_fills[ri][ci]; color = WHITE
            al = first_align if ci == 0 else PP_ALIGN.CENTER
            style_cell(tbl.cell(ri + 1, ci), val, fill=fill, color=color, align=al, size=size,
                       bold=(ci == 0 and val == "TOTAL"))
    return gf

def chunks(items, per):
    return [items[i:i + per] for i in range(0, len(items), per)]

def add_bar_chart(sl, categories, series, *, top=TOP_DFLT, horizontal=True, pct=False,
                  series_colors=None, highlight_negative=False):
    """Barres horizontales (lisibles pour beaucoup d'entités), avec étiquettes de valeurs.
       highlight_negative : colore en rouge foncé les barres dont la valeur est < 0."""
    cd = CategoryChartData()
    cd.categories = categories
    for nm, vals in series:
        cd.add_series(nm, vals)
    ctype = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = sl.shapes.add_chart(ctype, Cm(0.5), top, EMU_W - Cm(1.0),
                             int(EMU_H - top - BOT_MARGIN), cd)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.TOP
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(12)
    plot = ch.plots[0]
    plot.gap_width = 50
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(10); dl.font.bold = True
    dl.number_format = '0.0"%"' if pct else '0'
    dl.number_format_is_linked = False
    NEG = RGBColor(0x7F, 0x1D, 0x1D)
    for si, (nm, vals) in enumerate(series):
        ser = ch.series[si]
        if series_colors:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = series_colors[si]
        if highlight_negative:
            for pi, v in enumerate(vals):
                if v is not None and v < 0:
                    pt = ser.points[pi]
                    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = NEG
    try:
        ch.value_axis.tick_labels.font.size = Pt(11)
        ch.category_axis.tick_labels.font.size = Pt(11)
        if pct:
            ch.value_axis.has_major_gridlines = True
    except Exception:
        pass
    return gf

def chart_slides(title, sub, entity_chunks, label_fn, series_fns, names, *, after, pct=False,
                 series_colors=None, highlight_negative=False):
    """Crée une diapo par paquet d'entités ; titre numéroté (i/N)."""
    N = len(entity_chunks)
    for i, chunk in enumerate(entity_chunks, 1):
        sl = add_slide()
        suff = f"  ({i}/{N})" if N > 1 else ""
        add_title(sl, title + suff, sub)
        cats = [label_fn(k) for k in chunk]
        series = [(nm, [fn(k) for k in chunk]) for nm, fn in series_fns]
        add_bar_chart(sl, cats, series, pct=pct, series_colors=series_colors,
                      highlight_negative=highlight_negative)
        remember(sl, after)

def table_slides(title, sub, header, items, row_fn, *, after, size, col_w, row_cm,
                 total_row=None, per=None):
    """Tableau paginé proprement : TOTAL en tête de 1ʳᵉ page, en-tête répété, numéro (i/N)."""
    cap = per or rows_per_page(row_cm)
    # La 1ʳᵉ page accueille la ligne TOTAL : on y réserve une place.
    items = list(items)
    pages, idx, first = [], 0, True
    while idx < len(items):
        take = cap - (1 if (first and total_row is not None) else 0)
        pages.append(items[idx:idx + take]); idx += take; first = False
    if not pages:
        pages = [[]]
    N = len(pages)
    for i, page in enumerate(pages, 1):
        rows, fills = [], []
        if i == 1 and total_row is not None:
            rows.append(total_row[0]); fills.append(total_row[1])
        for k in page:
            r, f = row_fn(k); rows.append(r); fills.append(f)
        sl = add_slide()
        suff = f"  ({i}/{N})" if N > 1 else ""
        add_title(sl, title + suff, sub)
        add_table(sl, header, rows, size=size, col_w=col_w, cell_fills=fills, row_cm=row_cm)
        remember(sl, after)

NEW = []  # (after_index_0based, slide_element)  -- collected for reordering

def remember(sl, after0):
    NEW.append((after0, sl))

# ------------------------------------------------------------
# DIAPO 10 — remplir la table existante (Kin Est, 2025 vs 2026)
# ------------------------------------------------------------
def count_cs(recs, field):
    seen = {}
    for r in recs:
        k = (r.get("_ZS"), r.get("_AS"))
        seen[k] = seen.get(k, 0) + num(r.get(field))
    return sum(1 for v in seen.values() if v > 0)

def pr(recs, p, rl):
    return f"{round(s(recs, p))} / {round(s(recs, rl))}"

def slide10_rows(recs):
    return [
        str(count_cs(recs, "seances_avancees_realisees") or "—"),
        str(count_cs(recs, "seances_mobiles_realisees") or "—"),
        pr(recs, "seances_fixes_prevues", "seances_fixes_realisees"),
        pr(recs, "seances_avancees_prevues", "seances_avancees_realisees"),
        pr(recs, "seances_mobiles_prevues", "seances_mobiles_realisees"),
    ]

col25 = slide10_rows(prev)
col26 = slide10_rows(cur)
sl10 = prs.slides[9]
tbl10 = None
for shp in sl10.shapes:
    if shp.has_table:
        tbl10 = shp.table; break
if tbl10 is not None:
    for i in range(5):
        for ci, txt in ((1, col25[i]), (2, col26[i])):
            cell = tbl10.cell(i + 1, ci)
            cell.text = txt
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(12); run.font.bold = True

# ------------------------------------------------------------
# DIAPO 11 — analyses par AS et ZS (Jan-Avr 2026)
# ------------------------------------------------------------
SUB = f"Antenne Kinshasa Est — {PER_LABEL_26}"

# 1. Complétude & promptitude — graphiques (AS paginé, ZS unique)
def comp_prompt(recs):
    n = s(recs, "n"); rap = s(recs, "rap")
    pw = sum(num(r.get("prompt")) * num(r.get("n")) for r in recs)
    comp = rap / n * 100 if n else 0
    prompt = pw / n if n else 0
    return round(comp, 1), round(prompt, 1)

CHART_PER = 12   # entités max par diapo-graphique → barres plus grandes et lisibles
CP_COLORS = [RGBColor(0x00, 0x93, 0xD5), RGBColor(0xF5, 0x9E, 0x0B)]
ZDSV_COLORS = [RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xEA, 0xB3, 0x08)]
AB_COLORS = [RGBColor(0xC8, 0x1E, 0x1E), RGBColor(0xF5, 0x9E, 0x0B)]

cp_series = [("Complétude %", lambda k, g: comp_prompt(g[k])[0]),
             ("Promptitude %", lambda k, g: comp_prompt(g[k])[1])]
# ZS
chart_slides("Complétude & Promptitude par Zone de Santé", SUB, [zs_keys], lambda k: k,
             [(nm, (lambda f: lambda k: f(k, g_zs))(fn)) for nm, fn in cp_series],
             None, after=10, pct=True, series_colors=CP_COLORS)
# AS
chart_slides("Complétude & Promptitude par Aire de Santé", SUB, chunks(as_keys, CHART_PER), as_label,
             [(nm, (lambda f: lambda k: f(k, g_as))(fn)) for nm, fn in cp_series],
             None, after=10, pct=True, series_colors=CP_COLORS)

# 2. CV admin tous antigènes — tableaux (ZS unique, AS paginé)
def cv_row(recs):
    t = targets(recs)
    out = []
    fills = []
    for ag in AG_ORDER:
        p = cv_pct(recs, ag, t)
        out.append("–" if p is None else f"{p:.0f}")
        fills.append(cv_color(p))
    return out, fills

CV_COLW = [2.4] + [1] * len(AG_ORDER)
def cv_total():
    c, fl = cv_row(cur); return (["TOTAL"] + c, [None] + fl)
def cv_zs_row(k):
    c, fl = cv_row(g_zs[k]); return ([k] + c, [None] + fl)
def cv_as_row(k):
    c, fl = cv_row(g_as[k]); return ([as_label(k)] + c, [None] + fl)

table_slides("Couverture vaccinale administrative par ZS — tous antigènes (%)", SUB,
             ["Zone de Santé"] + AG_ORDER, zs_keys, cv_zs_row, after=10, size=10,
             col_w=CV_COLW, row_cm=0.95, total_row=cv_total())
table_slides("Couverture vaccinale administrative par AS — tous antigènes (%)", SUB,
             ["Aire de Santé"] + AG_ORDER, as_keys, cv_as_row, after=10, size=9,
             col_w=CV_COLW, row_cm=0.72, total_row=cv_total())

# 3. Enfants ZD & SV — graphiques
zdsv_series = [("Enfants ZD", lambda k, g: round(zd(g[k]))),
               ("Enfants SV", lambda k, g: round(sv(g[k])))]
chart_slides("Enfants Zéro-Dose (ZD) et Sous-Vaccinés (SV) par ZS", SUB, [zs_keys], lambda k: k,
             [(nm, (lambda f: lambda k: f(k, g_zs))(fn)) for nm, fn in zdsv_series], None, after=10,
             series_colors=ZDSV_COLORS)
chart_slides("Enfants Zéro-Dose (ZD) et Sous-Vaccinés (SV) par AS", SUB, chunks(as_keys, CHART_PER), as_label,
             [(nm, (lambda f: lambda k: f(k, g_as))(fn)) for nm, fn in zdsv_series], None, after=10,
             series_colors=ZDSV_COLORS)

# 4. Catégorisation ACZ — tableaux (ZS unique, AS paginé)
def cat_row(recs):
    t = targets(recs)
    d1 = s(recs, "DTC1_0_11"); d3 = s(recs, "DTC3_0_11")
    cv1 = d1 / t["ns"] * 100 if t["ns"] > 0 else None
    ta = abandon(d1, d3)
    cat = categorize(cv1, ta)
    return ["–" if cv1 is None else f"{cv1:.1f}%",
            "–" if ta is None else f"{ta:.1f}%", cat], cat, ta

CAT_NOTE = (SUB + " — Cat.1: P1≥90 & TA≤10 · Cat.2: P1≥90 & TA>10 · "
            "Cat.3: P1<90 & TA≤10 · Cat.4: P1<90 & TA>10")
CAT_COLW = [3.4, 2, 2.6, 2]
def _cat_fill(ta, cat):
    return [None, None, ab_color(ta), cat_color(cat)]
def cat_total():
    cells, cat, ta = cat_row(cur); return (["TOTAL"] + cells, _cat_fill(ta, cat))
def cat_zs_r(k):
    cells, cat, ta = cat_row(g_zs[k]); return ([k] + cells, _cat_fill(ta, cat))
def cat_as_r(k):
    cells, cat, ta = cat_row(g_as[k]); return ([as_label(k)] + cells, _cat_fill(ta, cat))

table_slides("Catégorisation ACZ par ZS", CAT_NOTE,
             ["Zone de Santé", "CV Penta1 (%)", "Taux abandon DTC1-DTC3 (%)", "Catégorie ACZ"],
             zs_keys, cat_zs_r, after=10, size=13, col_w=CAT_COLW, row_cm=1.05, total_row=cat_total())
table_slides("Catégorisation ACZ par AS", CAT_NOTE,
             ["Aire de Santé", "CV Penta1 (%)", "Taux abandon DTC1-DTC3 (%)", "Catégorie ACZ"],
             as_keys, cat_as_r, after=10, size=12, col_w=CAT_COLW, row_cm=0.9, total_row=cat_total())

# 5. Disponibilité des intrants (BCG, DTC, VAR, VPO) — tableaux
DISPO_V = [("BCG", "BCG"), ("DTC (Penta)", "DTC"), ("VAR", "VAR"), ("VPO", "VPO")]
def dispo_row(recs):
    cells, fills = [], []
    for _, pfx in DISPO_V:
        p = dispo(recs, pfx)
        cells.append("–" if p is None else f"{p:.1f}%"); fills.append(disp_color(p))
    return cells, fills

DISPO_COLW = [3.4, 1.5, 1.5, 1.5, 1.5]
def dispo_total():
    c, fl = dispo_row(cur); return (["TOTAL"] + c, [None] + fl)
table_slides("Taux de disponibilité des intrants par ZS (BCG, DTC, VAR, VPO)", SUB,
             ["Zone de Santé"] + [l for l, _ in DISPO_V], zs_keys,
             lambda k: (["%s" % k] + dispo_row(g_zs[k])[0], [None] + dispo_row(g_zs[k])[1]),
             after=10, size=13, col_w=DISPO_COLW, row_cm=1.05, total_row=dispo_total())
table_slides("Taux de disponibilité des intrants par AS (BCG, DTC, VAR, VPO)", SUB,
             ["Aire de Santé"] + [l for l, _ in DISPO_V], as_keys,
             lambda k: ([as_label(k)] + dispo_row(g_as[k])[0], [None] + dispo_row(g_as[k])[1]),
             after=10, size=12, col_w=DISPO_COLW, row_cm=0.9, total_row=dispo_total())

# 6. Taux de perte (VPO, DTC, VAR, VAA, Td) — tableaux
PERTE_V = [("VPO", "VPO"), ("DTC (Penta)", "DTC"), ("VAR", "VAR"), ("VAA", "VAA"), ("Td", "Td")]
def perte_row(recs):
    cells, fills = [], []
    for _, pfx in PERTE_V:
        p = perte(recs, pfx)
        cells.append("–" if p is None else f"{p:.1f}%"); fills.append(perte_color(p))
    return cells, fills

PERTE_COLW = [3.4, 1.4, 1.4, 1.4, 1.4, 1.4]
def perte_total():
    c, fl = perte_row(cur); return (["TOTAL"] + c, [None] + fl)
table_slides("Taux de perte par ZS et antigène traceur (VPO, DTC, VAR, VAA, Td)", SUB,
             ["Zone de Santé"] + [l for l, _ in PERTE_V], zs_keys,
             lambda k: (["%s" % k] + perte_row(g_zs[k])[0], [None] + perte_row(g_zs[k])[1]),
             after=10, size=13, col_w=PERTE_COLW, row_cm=1.05, total_row=perte_total())
table_slides("Taux de perte par AS et antigène traceur (VPO, DTC, VAR, VAA, Td)", SUB,
             ["Aire de Santé"] + [l for l, _ in PERTE_V], as_keys,
             lambda k: ([as_label(k)] + perte_row(g_as[k])[0], [None] + perte_row(g_as[k])[1]),
             after=10, size=12, col_w=PERTE_COLW, row_cm=0.9, total_row=perte_total())

# ------------------------------------------------------------
# DIAPO 12 — performance communication
# ------------------------------------------------------------
# slide12 original index : 11 (0-based). Find current index after additions? We append then reorder.
S12_AFTER = 11

# 1. Enfants perdus de vue identifiés & récupérés PAR LES RECO (hors BCU)
#    Indicateurs DHIS2 disponibles « par les RECO » : Penta1 + Penta3 (= DTC), tous âges.
PDV_NOTE = (SUB + " — Perdus de vue identifiés et récupérés par les RECO (DTC = Penta1 + Penta3, "
            "tous âges). Source DHIS2 ; indicateur « par les RECO » non collecté pour la VAR.")
PDV_HDR = ["{}", "Identifiés (RECO)", "Récupérés (RECO)", "% récupérés"]
PDV_COLW = [3.4, 2, 2, 1.6]
def _pdv(e):
    i, rc = e["id"], e["rec"]
    return f"{round(i)}", f"{round(rc)}", (f"{rc / i * 100:.0f}%" if i > 0 else "–")
def pdv_total():
    tid = sum(pdv_zs.get(k, {"id": 0})["id"] for k in zs_keys)
    trec = sum(pdv_zs.get(k, {"rec": 0})["rec"] for k in zs_keys)
    pct = f"{trec / tid * 100:.0f}%" if tid > 0 else "–"
    return (["TOTAL", f"{round(tid)}", f"{round(trec)}", pct], [None, None, None, None])
def pdv_zs_r(k):
    a, b, c = _pdv(pdv_zs.get(k, {"id": 0, "rec": 0})); return ([k, a, b, c], [None] * 4)
def pdv_as_r(k):
    a, b, c = _pdv(pdv_as.get(k, {"id": 0, "rec": 0})); return ([as_label(k), a, b, c], [None] * 4)

table_slides("Enfants perdus de vue identifiés et récupérés par les RECO — par ZS", PDV_NOTE,
             ["Zone de Santé"] + PDV_HDR[1:], zs_keys, pdv_zs_r, after=S12_AFTER,
             size=13, col_w=PDV_COLW, row_cm=1.05, total_row=pdv_total())
table_slides("Enfants perdus de vue identifiés et récupérés par les RECO — par AS", PDV_NOTE,
             ["Aire de Santé"] + PDV_HDR[1:], as_keys, pdv_as_r, after=S12_AFTER,
             size=12, col_w=PDV_COLW, row_cm=0.9, total_row=pdv_total())

# 2. Taux d'abandon BCG-VAR1 & DTC1-DTC3 — graphiques (ZS unique, AS paginé)
def ab_pair(recs):
    return (round(abandon(s(recs, "DTC1_0_11"), s(recs, "DTC3_0_11")) or 0, 1),
            round(abandon(s(recs, "BCG_0_11"), s(recs, "VAR1_0_11")) or 0, 1))

ab_series = [("Abandon DTC1-DTC3 %", lambda k, g: ab_pair(g[k])[0]),
             ("Abandon BCG-VAR1 %", lambda k, g: ab_pair(g[k])[1])]
chart_slides("Taux d'abandon DTC1-DTC3 et BCG-VAR1(RR1) par ZS", SUB, [zs_keys], lambda k: k,
             [(nm, (lambda f: lambda k: f(k, g_zs))(fn)) for nm, fn in ab_series], None, after=S12_AFTER,
             pct=True, series_colors=AB_COLORS, highlight_negative=True)
chart_slides("Taux d'abandon DTC1-DTC3 et BCG-VAR1(RR1) par AS", SUB, chunks(as_keys, CHART_PER), as_label,
             [(nm, (lambda f: lambda k: f(k, g_as))(fn)) for nm, fn in ab_series], None, after=S12_AFTER,
             pct=True, series_colors=AB_COLORS, highlight_negative=True)

# ------------------------------------------------------------
#  RÉORDONNANCEMENT : insérer les nouvelles diapos après 11 / 12
# ------------------------------------------------------------
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)              # ordre courant (27 base + nouvelles appendues)
base = ids[:27]                   # 27 diapos d'origine
# map slide element -> sldId
def sldid_of(slide):
    rId = None
    for sid in ids:
        if sid.get(qn('r:id')) == slide.part.relate_to.__self__ if False else False:
            pass
    return None

# Construire ordre : pour chaque diapo de base, l'ajouter ; après index 10 (diapo11)
# insérer les NEW after0==10 ; après index 11 (diapo12) insérer NEW after0==11.
# On retrouve l'élément sldId de chaque slide via son position dans prs.slides.
slide_to_sldId = {}
for sid in ids:
    slide_to_sldId[sid] = sid
# Build position list of new slides grouped
new11 = [sl for (a, sl) in NEW if a == 10]
new12 = [sl for (a, sl) in NEW if a == 11]

def sldId_for_slide(slide):
    target_rid = None
    part = slide.part
    # find rId in presentation rels pointing to this part
    for rid, rel in prs.part.rels.items():
        if rel._target is part:
            target_rid = rid; break
    for sid in list(sldIdLst):
        if sid.get(qn('r:id')) == target_rid:
            return sid
    return None

order = []
for i, sid in enumerate(base):
    order.append(sid)
    if i == 10:   # diapo 11
        for sl in new11:
            order.append(sldId_for_slide(sl))
    if i == 11:   # diapo 12
        for sl in new12:
            order.append(sldId_for_slide(sl))

# Réécrire sldIdLst dans le nouvel ordre
for sid in list(sldIdLst):
    sldIdLst.remove(sid)
for sid in order:
    sldIdLst.append(sid)

prs.save(OUT)
print("OK ->", OUT)
print("Diapos totales:", len(prs.slides._sldIdLst))
print("Nouvelles diapo 11:", len(new11), "| Nouvelles diapo 12:", len(new12))
